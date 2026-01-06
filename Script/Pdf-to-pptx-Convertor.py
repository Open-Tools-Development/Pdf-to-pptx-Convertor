import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import filedialog, messagebox

def convert_pdf_to_pptx(pdf_path, output_folder):
    pptx_filename = os.path.splitext(os.path.basename(pdf_path))[0] + ".pptx"
    output_path = os.path.join(output_folder, pptx_filename)

    images = convert_from_path(pdf_path, dpi=200)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for i, img in enumerate(images):
        temp_img = f"temp_slide_{i}.png"
        img.save(temp_img, "PNG")
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
        os.remove(temp_img)

    prs.save(output_path)
    print(f"Saved: {output_path}")

def main():
    root = tk.Tk()
    root.withdraw()

    choice = messagebox.askyesno("PDF to PPTX", "Do you want to convert a **folder** of PDFs?\n(Yes = Folder, No = Single File)")

    if choice:
        pdf_folder = filedialog.askdirectory(title="Select Folder Containing PDF Files")
        if not pdf_folder:
            print("No folder selected.")
            return
        pdf_files = [os.path.join(pdf_folder, f) for f in os.listdir(pdf_folder) if f.lower().endswith(".pdf")]
    else:
        pdf_file = filedialog.askopenfilename(title="Select a PDF File", filetypes=[("PDF files", "*.pdf")])
        if not pdf_file:
            print("No file selected.")
            return
        pdf_files = [pdf_file]

    output_folder = filedialog.askdirectory(title="Select Destination Folder")
    if not output_folder:
        print("No output folder selected.")
        return

    for pdf in pdf_files:
        try:
            convert_pdf_to_pptx(pdf, output_folder)
        except Exception as e:
            print(f"Failed to convert {pdf}: {e}")

if __name__ == "__main__":
    main()
